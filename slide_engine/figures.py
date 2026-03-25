"""
Scientific figure builder for PowerPoint.

Produces multi-panel figures with independent, movable, editable objects:
- Panel grid layout (A, B, C, D...) with automatic labeling
- Chart panels (matplotlib → embedded image)
- Diagram panels (sem_engine → native shapes)
- Schematic panels (icon + connector compositions)
- Annotation layer (arrows, brackets, callouts)
- Scientific styling presets (Nature, Cell, general)

Every object is a native PowerPoint shape — independently movable and
editable when opened in Microsoft PowerPoint.
"""

from __future__ import annotations

import math
from pathlib import Path
from dataclasses import dataclass, field
from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ── Style presets ────────────────────────────────────────────────────────

@dataclass
class FigureStyle:
    name: str
    font: str = "Arial"
    panel_label_size: int = 16        # Pt — the A, B, C labels
    panel_label_bold: bool = True
    title_size: int = 11
    axis_label_size: int = 9
    tick_size: int = 8
    annotation_size: int = 8
    caption_size: int = 8
    line_width: float = 1.0           # Pt
    color_primary: str = "#333333"
    color_secondary: str = "#666666"
    color_accent: str = "#2E86C1"
    color_background: str = "#FFFFFF"
    panel_border: bool = False
    panel_padding: float = 0.15       # inches inside panel


STYLE_NATURE = FigureStyle(
    name="nature", font="Arial",
    panel_label_size=16, title_size=10, axis_label_size=8,
    color_primary="#000000", color_accent="#E31A1C",
)

STYLE_CELL = FigureStyle(
    name="cell", font="Helvetica",
    panel_label_size=14, title_size=10, axis_label_size=8,
    color_primary="#333333", color_accent="#1F78B4",
)

STYLE_GENERAL = FigureStyle(
    name="general", font="Calibri",
    panel_label_size=14, title_size=11, axis_label_size=9,
    color_primary="#2C3E50", color_accent="#2E86C1",
)

STYLE_DARK = FigureStyle(
    name="dark", font="Calibri",
    panel_label_size=14, title_size=11, axis_label_size=9,
    color_primary="#E0E0E0", color_secondary="#90A4AE",
    color_accent="#00A0DC", color_background="#0D1B2A",
    panel_border=True,
)

FIGURE_STYLES = {
    "nature": STYLE_NATURE,
    "cell": STYLE_CELL,
    "general": STYLE_GENERAL,
    "dark": STYLE_DARK,
}


# ── Panel definitions ────────────────────────────────────────────────────

@dataclass
class Panel:
    """A single panel in a multi-panel figure."""
    label: str = ""                   # "A", "B", etc.
    title: str = ""                   # Panel title/caption
    panel_type: str = "empty"         # empty, chart_image, diagram, schematic, image
    row: int = 0
    col: int = 0
    row_span: int = 1
    col_span: int = 1
    # For chart_image type
    image_path: str = ""
    # For diagram type (renders via sem_engine)
    flow_spec: dict | None = None
    # For schematic type (icons + connectors)
    elements: list[dict] = field(default_factory=list)
    # Annotations on this panel
    annotations: list[dict] = field(default_factory=list)


@dataclass
class Figure:
    """A multi-panel scientific figure."""
    title: str = ""
    rows: int = 2
    cols: int = 2
    panels: list[Panel] = field(default_factory=list)
    style: str = "general"
    # Overall figure dimensions on the slide (inches)
    left: float = 0.5
    top: float = 1.4
    width: float = 12.3
    height: float = 5.5


# ── Renderer ─────────────────────────────────────────────────────────────

def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class FigureRenderer:
    """Renders a Figure onto a PowerPoint slide as independent editable objects."""

    def __init__(self, style: FigureStyle | None = None):
        self.style = style or STYLE_GENERAL

    def render(self, slide, figure: Figure):
        """Render the complete figure onto a slide."""
        self.style = FIGURE_STYLES.get(figure.style, STYLE_GENERAL)

        # Calculate panel grid positions
        gap_x = 0.25  # inches between panels
        gap_y = 0.3
        panel_w = (figure.width - gap_x * (figure.cols - 1)) / figure.cols
        panel_h = (figure.height - gap_y * (figure.rows - 1)) / figure.rows

        for panel in figure.panels:
            # Calculate panel position
            px = figure.left + panel.col * (panel_w + gap_x)
            py = figure.top + panel.row * (panel_h + gap_y)
            pw = panel_w * panel.col_span + gap_x * (panel.col_span - 1)
            ph = panel_h * panel.row_span + gap_y * (panel.row_span - 1)

            # Draw panel background/border
            if self.style.panel_border:
                border = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(px), Inches(py), Inches(pw), Inches(ph)
                )
                border.fill.background()
                border.line.color.rgb = _rgb(self.style.color_secondary)
                border.line.width = Pt(0.5)

            # Panel label (A, B, C...)
            if panel.label:
                lbl = slide.shapes.add_textbox(
                    Inches(px + 0.05), Inches(py + 0.02),
                    Inches(0.4), Inches(0.3)
                )
                p = lbl.text_frame.paragraphs[0]
                p.text = panel.label
                p.font.size = Pt(self.style.panel_label_size)
                p.font.bold = self.style.panel_label_bold
                p.font.name = self.style.font
                p.font.color.rgb = _rgb(self.style.color_primary)

            # Panel title
            if panel.title:
                ttl = slide.shapes.add_textbox(
                    Inches(px + 0.4), Inches(py + 0.05),
                    Inches(pw - 0.5), Inches(0.25)
                )
                p = ttl.text_frame.paragraphs[0]
                p.text = panel.title
                p.font.size = Pt(self.style.title_size)
                p.font.bold = True
                p.font.name = self.style.font
                p.font.color.rgb = _rgb(self.style.color_primary)

            # Content area (inside panel, below label)
            content_top = py + 0.35
            content_height = ph - 0.4
            content_left = px + self.style.panel_padding
            content_width = pw - 2 * self.style.panel_padding

            # Render panel content based on type
            if panel.panel_type == "chart_image" and panel.image_path:
                self._render_chart_image(slide, panel, content_left, content_top, content_width, content_height)
            elif panel.panel_type == "diagram" and panel.flow_spec:
                self._render_diagram(slide, panel, content_left, content_top, content_width, content_height)
            elif panel.panel_type == "schematic":
                self._render_schematic(slide, panel, content_left, content_top, content_width, content_height)
            elif panel.panel_type == "image" and panel.image_path:
                self._render_image(slide, panel, content_left, content_top, content_width, content_height)

            # Annotations
            for ann in panel.annotations:
                self._render_annotation(slide, ann, px, py, pw, ph)

    def _render_chart_image(self, slide, panel, left, top, width, height):
        """Embed a matplotlib-generated chart image."""
        img_path = Path(panel.image_path)
        if img_path.exists():
            slide.shapes.add_picture(
                str(img_path),
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )

    def _render_image(self, slide, panel, left, top, width, height):
        """Embed a static image (icon, photo, etc.)."""
        img_path = Path(panel.image_path)
        if img_path.exists():
            slide.shapes.add_picture(
                str(img_path),
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )

    def _render_diagram(self, slide, panel, left, top, width, height):
        """Render a process flow diagram as native shapes within the panel."""
        try:
            from sem_engine.schema import ProcessModel
            from sem_engine.renderers.pptx_native import PptxNativeFlowRenderer

            model = ProcessModel.from_dict(panel.flow_spec)
            dark = self.style.name == "dark"
            renderer = PptxNativeFlowRenderer(
                dark_mode=dark, node_w=1.2, node_h=0.45,
                font_size=8, label_font_size=7, group_label_size=8
            )
            renderer.render_on_slide(slide, model, left=left, top=top, width=width, height=height)
        except Exception as e:
            # Fallback: placeholder text
            txbox = slide.shapes.add_textbox(
                Inches(left), Inches(top + height/3),
                Inches(width), Inches(0.3)
            )
            p = txbox.text_frame.paragraphs[0]
            p.text = f"[Diagram: {e}]"
            p.font.size = Pt(8)
            p.font.color.rgb = _rgb(self.style.color_secondary)
            p.alignment = PP_ALIGN.CENTER

    def _render_schematic(self, slide, panel, left, top, width, height):
        """Render a schematic composed of icons and connectors."""
        for elem in panel.elements:
            etype = elem.get("type", "box")
            ex = left + elem.get("x", 0) * width
            ey = top + elem.get("y", 0) * height
            ew = elem.get("w", 0.15) * width
            eh = elem.get("h", 0.15) * height
            color = elem.get("color", self.style.color_accent)
            label = elem.get("label", "")

            if etype == "box":
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(ex), Inches(ey), Inches(ew), Inches(eh)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = _rgb(color)
                shape.line.fill.background()
                try:
                    shape.adjustments[0] = 0.2
                except Exception:
                    pass

                if label:
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = label
                    p.font.size = Pt(self.style.axis_label_size)
                    p.font.bold = True
                    p.font.name = self.style.font
                    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    p.alignment = PP_ALIGN.CENTER

            elif etype == "circle":
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(ex), Inches(ey), Inches(ew), Inches(eh)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = _rgb(color)
                shape.line.fill.background()

                if label:
                    tf = shape.text_frame
                    p = tf.paragraphs[0]
                    p.text = label
                    p.font.size = Pt(self.style.axis_label_size - 1)
                    p.font.bold = True
                    p.font.name = self.style.font
                    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    p.alignment = PP_ALIGN.CENTER

            elif etype == "arrow":
                # Arrow from (x,y) to (x2,y2) in relative coordinates
                x2 = left + elem.get("x2", 0.5) * width
                y2 = top + elem.get("y2", 0.5) * height
                connector = slide.shapes.add_connector(
                    1, Inches(ex + ew/2), Inches(ey + eh/2),
                    Inches(x2), Inches(y2)
                )
                connector.line.color.rgb = _rgb(color)
                connector.line.width = Pt(self.style.line_width)

            elif etype == "icon" and elem.get("image"):
                img_path = Path(elem["image"])
                if img_path.exists():
                    slide.shapes.add_picture(
                        str(img_path),
                        Inches(ex), Inches(ey), Inches(ew), Inches(eh)
                    )
                if label:
                    lbl = slide.shapes.add_textbox(
                        Inches(ex), Inches(ey + eh + 0.02),
                        Inches(ew), Inches(0.2)
                    )
                    p = lbl.text_frame.paragraphs[0]
                    p.text = label
                    p.font.size = Pt(self.style.tick_size)
                    p.font.name = self.style.font
                    p.font.color.rgb = _rgb(self.style.color_primary)
                    p.alignment = PP_ALIGN.CENTER

            elif etype == "text":
                txbox = slide.shapes.add_textbox(
                    Inches(ex), Inches(ey), Inches(ew), Inches(eh)
                )
                p = txbox.text_frame.paragraphs[0]
                p.text = label
                p.font.size = Pt(elem.get("size", self.style.annotation_size))
                p.font.bold = elem.get("bold", False)
                p.font.name = self.style.font
                p.font.color.rgb = _rgb(elem.get("text_color", self.style.color_primary))
                p.alignment = PP_ALIGN.CENTER if elem.get("center", False) else PP_ALIGN.LEFT

    def _render_annotation(self, slide, ann: dict, px, py, pw, ph):
        """Render an annotation (arrow, bracket, callout) on a panel."""
        ann_type = ann.get("type", "callout")
        text = ann.get("text", "")

        # Relative positions within panel
        x = px + ann.get("x", 0.5) * pw
        y = py + ann.get("y", 0.5) * ph

        if ann_type == "callout":
            txbox = slide.shapes.add_textbox(
                Inches(x), Inches(y),
                Inches(ann.get("w", 1.5)), Inches(0.25)
            )
            p = txbox.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(self.style.annotation_size)
            p.font.name = self.style.font
            p.font.color.rgb = _rgb(ann.get("color", self.style.color_accent))
            p.font.italic = True

        elif ann_type == "arrow":
            x2 = px + ann.get("x2", 0.6) * pw
            y2 = py + ann.get("y2", 0.6) * ph
            connector = slide.shapes.add_connector(
                1, Inches(x), Inches(y), Inches(x2), Inches(y2)
            )
            connector.line.color.rgb = _rgb(ann.get("color", self.style.color_accent))
            connector.line.width = Pt(1.5)
