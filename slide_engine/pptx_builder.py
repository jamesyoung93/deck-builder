"""
PowerPoint rendering engine using python-pptx.

Translates a Deck schema into a .pptx file with consulting-quality formatting.
Design patterns derived from consulting, consulting, and consulting reference materials.
"""

from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .schema import Deck, Slide, SlideType, BulletPoint
from .styles import Style, get_style, Grid
from .icons import get_icon_path, auto_icon


class PptxBuilder:

    def __init__(self, style: Style | None = None):
        self.style = style or get_style("neutral")
        self._init_style()

    def _init_style(self):
        self.grid = self.style.grid
        self.colors = self.style.colors
        self.typo = self.style.typography

    def build(self, deck: Deck, output_path: str | Path) -> Path:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self.style = get_style(deck.style.value)
        self._init_style()

        prs = Presentation()
        prs.slide_width = self.grid.slide_width
        prs.slide_height = self.grid.slide_height

        for i, spec in enumerate(deck.slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._apply_bg(slide)
            self._render(slide, spec, i, len(deck.slides))

        prs.save(str(output_path))
        return output_path

    # ── Background ───────────────────────────────────────────────────────

    def _apply_bg(self, slide, color_key: str = "background"):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors.rgb(color_key)

    def _rgb(self, hex_str: str) -> RGBColor:
        return self.colors.rgb_hex(hex_str)

    # ── Chrome: title, accent line, footer ───────────────────────────────

    def _add_title(self, slide, title: str, subtitle: str = ""):
        # Title text
        txbox = slide.shapes.add_textbox(
            self.grid.title_left, self.grid.title_top,
            self.grid.title_width, self.grid.title_height
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(self.typo.title_size)
        p.font.bold = self.typo.title_bold
        p.font.name = self.typo.title_font
        p.font.color.rgb = self.colors.rgb("title_text")

        # Subtitle under title
        if subtitle:
            sub = tf.add_paragraph()
            sub.text = subtitle
            sub.font.size = Pt(self.typo.subtitle_size - 2)
            sub.font.name = self.typo.body_font
            sub.font.color.rgb = self.colors.rgb("subtitle_text")
            sub.space_before = Pt(4)

        # Accent line (only for light-background styles)
        if self.style.title_accent_line and self.style.title_accent_color:
            line_top = self.grid.title_top + self.grid.title_height + Inches(0.05)
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                self.grid.title_left, line_top,
                self.grid.title_width, Inches(0.025)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = self._rgb(self.style.title_accent_color)
            line.line.fill.background()

    def _add_footer(self, slide, page: int, total: int, source: str = ""):
        # Thin separator line above footer
        self._rect(slide, self.grid.margin_left, self.grid.footer_top - Inches(0.08),
                   self.grid.content_width, Inches(0.008),
                   self.colors.border)

        # Page number
        pg = slide.shapes.add_textbox(
            self.grid.slide_width - Inches(0.8), self.grid.footer_top,
            Inches(0.5), self.grid.footer_height
        )
        p = pg.text_frame.paragraphs[0]
        p.text = str(page)
        p.font.size = Pt(self.typo.footer_size)
        p.font.color.rgb = self.colors.rgb("footer_text")
        p.font.name = self.typo.body_font
        p.alignment = PP_ALIGN.RIGHT

        # Source citation
        if source:
            src = slide.shapes.add_textbox(
                self.grid.margin_left, self.grid.footer_top,
                Inches(9), self.grid.footer_height
            )
            p = src.text_frame.paragraphs[0]
            p.text = f"Source: {source}"
            p.font.size = Pt(self.typo.source_size)
            p.font.color.rgb = self.colors.rgb("footer_text")
            p.font.name = self.typo.body_font
            p.font.italic = True

    # ── Helpers ──────────────────────────────────────────────────────────

    def _text(self, slide, text, left, top, width, height,
              size=None, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size or self.typo.body_size)
        p.font.bold = bold
        p.font.name = self.typo.body_font
        p.font.color.rgb = color or self.colors.rgb("body_text")
        p.font.italic = italic
        p.alignment = align
        return txbox

    def _card(self, slide, left, top, width, height, fill_hex: str = None):
        """Add a rounded rectangle card (subtle container)."""
        if fill_hex is None:
            # Slightly lighter than background for dark mode, slightly darker for light
            if self.style.dark_mode:
                fill_hex = "#142538"
            else:
                fill_hex = "#F5F7FA"
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill_hex)
        # Subtle border
        shape.line.color.rgb = self.colors.rgb("border")
        shape.line.width = Pt(0.5)
        # Adjust corner radius
        shape.adjustments[0] = 0.03
        return shape

    def _rect(self, slide, left, top, width, height, fill_hex: str):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill_hex)
        shape.line.fill.background()
        return shape

    def _add_icon(self, slide, icon_name: str, left, top, size=Inches(0.35),
                  color: str = "white"):
        """Add a small icon image to the slide."""
        icon_path = get_icon_path(icon_name, color)
        if icon_path and icon_path.exists():
            slide.shapes.add_picture(str(icon_path), left, top, size, size)
            return True
        return False

    def _add_rich_bullets(self, slide, bullets: list[BulletPoint],
                          left, top, width, height, with_icons: bool = True):
        """Bullets with blue bold lead-in and auto-detected icons (consulting style)."""
        lead_color = self._rgb(self.colors.primary) if self.style.dark_mode else self.colors.rgb("body_text")
        icon_color = "blue" if self.style.dark_mode else "navy"
        icon_size = Inches(0.3)
        icon_indent = Inches(0.45)

        y = top
        row_height = Inches(0.65)

        for b in bullets:
            # Auto-detect icon
            if with_icons:
                icon_name = auto_icon(b.lead + " " + b.detail)
                if icon_name:
                    self._add_icon(slide, icon_name, left, y + Inches(0.02), icon_size, icon_color)

            # Text with bold lead + detail
            text_left = left + (icon_indent if with_icons else 0)
            text_width = width - (icon_indent if with_icons else 0)

            txbox = slide.shapes.add_textbox(text_left, y, text_width, row_height)
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]

            run = p.add_run()
            run.text = b.lead + (". " if b.detail else "")
            run.font.bold = True
            run.font.size = Pt(self.typo.bullet_size)
            run.font.name = self.typo.body_font
            run.font.color.rgb = lead_color

            if b.detail:
                run2 = p.add_run()
                run2.text = b.detail
                run2.font.bold = False
                run2.font.size = Pt(self.typo.bullet_size)
                run2.font.name = self.typo.body_font
                run2.font.color.rgb = self.colors.rgb("body_text")

            y += row_height

    # ── Dispatch ─────────────────────────────────────────────────────────

    def _render(self, slide, spec: Slide, idx: int, total: int):
        dispatch = {
            SlideType.COVER: self._cover,
            SlideType.AGENDA: self._agenda,
            SlideType.SECTION_DIVIDER: self._divider,
            SlideType.EXECUTIVE_SUMMARY: self._exec_summary,
            SlideType.ACTION_BULLETS: self._action_bullets,
            SlideType.TWO_COLUMN: self._two_col,
            SlideType.THREE_COLUMN: self._three_col,
            SlideType.DATA_CALLOUT: self._data_callout,
            SlideType.BAR_CHART: self._bar_chart,
            SlideType.PROCESS_FLOW: self._process_flow,
            SlideType.FRAMEWORK: self._framework,
            SlideType.TIMELINE: self._timeline,
            SlideType.QUOTE: self._quote,
            SlideType.CLOSING: self._closing,
        }
        dispatch.get(spec.type, self._action_bullets)(slide, spec, idx, total)

    # ── COVER ────────────────────────────────────────────────────────────

    def _cover(self, slide, s: Slide, idx: int, total: int):
        # Dark cover with accent stripe
        if self.style.dark_mode:
            # Thin accent bar at top
            self._rect(slide, 0, 0, self.grid.slide_width, Inches(0.06), self.colors.primary)
        else:
            # Full primary color background
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self._rgb(self.colors.primary)

        title_color = RGBColor(0xFF, 0xFF, 0xFF) if not self.style.dark_mode else self.colors.rgb("title_text")

        # Title
        self._text(slide, s.title,
                   Inches(1.2), Inches(2.0), Inches(10), Inches(2),
                   size=34, bold=True, color=title_color)

        # Thin line separator
        sep_color = self.colors.primary if self.style.dark_mode else "#FFFFFF"
        self._rect(slide, Inches(1.2), Inches(3.8), Inches(3), Inches(0.03), sep_color)

        # Subtitle
        if s.subtitle:
            sub_color = self.colors.rgb("subtitle_text") if self.style.dark_mode else self._rgb("#CCCCCC")
            self._text(slide, s.subtitle,
                       Inches(1.2), Inches(4.1), Inches(8), Inches(0.8),
                       size=16, color=sub_color)

        # Date
        date_text = s.body or s.notes or ""
        if date_text:
            dt_color = self.colors.rgb("footer_text") if self.style.dark_mode else self._rgb("#AAAAAA")
            self._text(slide, date_text,
                       Inches(1.2), Inches(5.2), Inches(5), Inches(0.4),
                       size=11, color=dt_color)

    # ── AGENDA ───────────────────────────────────────────────────────────

    def _agenda(self, slide, s: Slide, idx: int, total: int):
        self._add_title(slide, s.title or "Agenda")
        self._add_footer(slide, idx + 1, total)

        items = s.agenda_items or [b.lead for b in s.bullets]
        y = self.grid.content_top + Inches(0.4)

        for i, item in enumerate(items):
            is_active = (i == s.current_section)
            circle_x = Inches(1.5)

            # Number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, circle_x, y, Inches(0.45), Inches(0.45)
            )
            if is_active:
                circle.fill.solid()
                circle.fill.fore_color.rgb = self._rgb(self.colors.primary)
            else:
                circle.fill.solid()
                c = "#1E3A5F" if self.style.dark_mode else "#E8E8E8"
                circle.fill.fore_color.rgb = self._rgb(c)
            circle.line.fill.background()

            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

            # Item text
            txt_color = self.colors.rgb("body_text") if is_active else self.colors.rgb("subtitle_text")
            self._text(slide, item,
                       Inches(2.2), y + Inches(0.05), Inches(8), Inches(0.4),
                       size=14 if is_active else 13, bold=is_active, color=txt_color)

            y += Inches(0.65)

    # ── SECTION DIVIDER ──────────────────────────────────────────────────

    def _divider(self, slide, s: Slide, idx: int, total: int):
        if not self.style.dark_mode:
            # Light styles: use primary color background
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self._rgb(self.colors.divider_bg)
            text_color = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            # Dark mode: keep dark bg, use accent color for text
            text_color = self._rgb(self.colors.primary)
            # Add subtle accent line
            self._rect(slide, Inches(1.2), Inches(3.0), Inches(2.5), Inches(0.04), self.colors.primary)

        self._text(slide, s.title,
                   Inches(1.2), Inches(3.2), Inches(10), Inches(2),
                   size=32, bold=True, color=text_color)

        self._add_footer(slide, idx + 1, total)

    # ── EXECUTIVE SUMMARY ────────────────────────────────────────────────

    def _exec_summary(self, slide, s: Slide, idx: int, total: int):
        self._add_title(slide, s.title)
        self._add_footer(slide, idx + 1, total, s.source)

        n = len(s.bullets)
        if n == 0:
            return

        # Calculate row height to fill available space evenly
        avail_h = self.grid.content_height - Inches(0.6)
        row_h_raw = avail_h / n
        row_h = min(row_h_raw, Inches(1.1))  # cap row height
        y = self.grid.content_top + Inches(0.3)

        lead_color = self._rgb(self.colors.primary) if self.style.dark_mode else self.colors.rgb("body_text")
        icon_color = "blue" if self.style.dark_mode else "navy"
        stripe_color = self.colors.primary if self.style.dark_mode else self.colors.border
        left_margin = self.grid.margin_left + Inches(0.2)

        for i, b in enumerate(s.bullets):
            # Left accent stripe
            self._rect(slide, left_margin, y + Inches(0.05),
                       Inches(0.04), row_h - Inches(0.15), stripe_color)

            # Icon
            icon_name = auto_icon(b.lead + " " + b.detail)
            icon_left = left_margin + Inches(0.2)
            if icon_name:
                self._add_icon(slide, icon_name, icon_left, y + Inches(0.08),
                              Inches(0.3), icon_color)

            # Text
            text_left = icon_left + Inches(0.42)
            text_w = self.grid.content_width - (text_left - self.grid.margin_left) - Inches(0.3)

            txbox = slide.shapes.add_textbox(text_left, y + Inches(0.05),
                                            text_w, row_h - Inches(0.1))
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]

            run = p.add_run()
            run.text = b.lead + ": " if b.detail else b.lead
            run.font.bold = True
            run.font.size = Pt(self.typo.bullet_size)
            run.font.name = self.typo.body_font
            run.font.color.rgb = lead_color

            if b.detail:
                run2 = p.add_run()
                run2.text = b.detail
                run2.font.bold = False
                run2.font.size = Pt(self.typo.bullet_size)
                run2.font.name = self.typo.body_font
                run2.font.color.rgb = self.colors.rgb("body_text")

            y += row_h

    # ── ACTION BULLETS ───────────────────────────────────────────────────

    def _action_bullets(self, slide, s: Slide, idx: int, total: int):
        self._add_title(slide, s.title)
        self._add_footer(slide, idx + 1, total, s.source)

        bullet_top = self.grid.content_top + Inches(0.2)

        if s.subtitle:
            self._text(slide, s.subtitle,
                       self.grid.margin_left + Inches(0.2), self.grid.content_top,
                       self.grid.content_width, Inches(0.4),
                       size=self.typo.subtitle_size, bold=True,
                       color=self._rgb(self.colors.primary))
            bullet_top += Inches(0.5)

        self._add_rich_bullets(slide, s.bullets,
                               self.grid.margin_left + Inches(0.4), bullet_top,
                               self.grid.content_width - Inches(0.8),
                               self.grid.content_height - Inches(1))

    # ── TWO COLUMN ───────────────────────────────────────────────────────

    def _two_col(self, slide, s: Slide, idx: int, total: int):
        self._add_title(slide, s.title)
        self._add_footer(slide, idx + 1, total, s.source)

        cols = s.columns[:2] if len(s.columns) >= 2 else s.columns
        gap = Inches(0.4)
        col_w = self.grid.col_width(2, gap)
        top = self.grid.content_top + Inches(0.2)

        for ci, col in enumerate(cols):
            left = self.grid.col_left(ci, 2, gap)

            # Card background
            card_h = self.grid.content_height - Inches(0.5)
            self._card(slide, left, top, col_w, card_h)

            # Column header
            self._text(slide, col.heading,
                       left + Inches(0.25), top + Inches(0.2), col_w - Inches(0.5), Inches(0.4),
                       size=self.typo.subtitle_size, bold=True,
                       color=self._rgb(self.colors.primary))

            # Separator line inside card
            self._rect(slide, left + Inches(0.25), top + Inches(0.6),
                       col_w - Inches(0.5), Inches(0.02), self.colors.border)

            # Bullets
            y = top + Inches(0.8)
            for bt in col.bullets:
                self._text(slide, bt,
                           left + Inches(0.35), y, col_w - Inches(0.7), Inches(0.35),
                           size=self.typo.bullet_size)
                y += Inches(0.38)

    # ── THREE COLUMN ─────────────────────────────────────────────────────

    def _three_col(self, slide, s: Slide, idx: int, total: int):
        self._add_title(slide, s.title)
        self._add_footer(slide, idx + 1, total, s.source)

        cols = s.columns[:3] if len(s.columns) >= 3 else s.columns
        gap = Inches(0.3)
        col_w = self.grid.col_width(3, gap)
        top = self.grid.content_top + Inches(0.2)
        icon_color = "white"

        for ci, col in enumerate(cols):
            left = self.grid.col_left(ci, 3, gap)

            # Card
            card_h = self.grid.content_height - Inches(0.5)
            self._card(slide, left, top, col_w, card_h)

            # Icon + number badge row
            badge_size = Inches(0.5)
            # Auto-detect icon for column heading
            col_icon = auto_icon(col.heading + " " + " ".join(col.bullets[:2]))
            icon_x = left + int(col_w / 2) - int(badge_size / 2)

            if col_icon:
                # Icon in a colored circle
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, icon_x, top + Inches(0.2), badge_size, badge_size
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self._rgb(self.colors.primary)
                circle.line.fill.background()
                self._add_icon(slide, col_icon, icon_x + Inches(0.1),
                              top + Inches(0.3), Inches(0.3), icon_color)
            else:
                # Number badge fallback
                badge = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, icon_x, top + Inches(0.2), badge_size, badge_size
                )
                badge.fill.solid()
                badge.fill.fore_color.rgb = self._rgb(self.colors.primary)
                badge.line.fill.background()
                p = badge.text_frame.paragraphs[0]
                p.text = str(ci + 1)
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.alignment = PP_ALIGN.CENTER

            # Heading
            self._text(slide, col.heading,
                       left + Inches(0.15), top + Inches(0.85),
                       col_w - Inches(0.3), Inches(0.4),
                       size=13, bold=True, align=PP_ALIGN.CENTER)

            # Separator
            self._rect(slide, left + Inches(0.3), top + Inches(1.3),
                       col_w - Inches(0.6), Inches(0.015), self.colors.border)

            # Bullets with small dot markers
            y = top + Inches(1.5)
            for bt in col.bullets:
                # Dot marker
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    left + Inches(0.25), y + Inches(0.1),
                    Inches(0.06), Inches(0.06)
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = self._rgb(self.colors.primary)
                dot.line.fill.background()

                self._text(slide, bt,
                           left + Inches(0.4), y, col_w - Inches(0.6), Inches(0.32),
                           size=self.typo.bullet_size - 1)
                y += Inches(0.35)

    # ── DATA CALLOUT ─────────────────────────────────────────────────────

    def _data_callout(self, slide, s: Slide, idx: int, total: int):
        self._add_title(slide, s.title)
        self._add_footer(slide, idx + 1, total, s.source)

        callouts = s.callouts
        n = max(len(callouts), 1)
        gap = Inches(0.25)
        col_w = self.grid.col_width(n, gap)
        top = self.grid.content_top + Inches(0.3)
        card_h = Inches(3.2)

        for ci, c in enumerate(callouts):
            left = self.grid.col_left(ci, n, gap)
            accent = c.color or self.colors.primary

            # Card container
            self._card(slide, left, top, col_w, card_h)

            # Accent bar at top of card
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, Inches(0.06)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = self._rgb(accent)
            bar.line.fill.background()
            try:
                bar.adjustments[0] = 0.5
            except Exception:
                pass

            # Large value
            self._text(slide, c.value,
                       left, top + Inches(0.3), col_w, Inches(1.2),
                       size=46, bold=True, color=self._rgb(accent),
                       align=PP_ALIGN.CENTER)

            # Label - bold
            self._text(slide, c.label,
                       left + Inches(0.15), top + Inches(1.55),
                       col_w - Inches(0.3), Inches(0.45),
                       size=11, bold=True, align=PP_ALIGN.CENTER)

            # Thin separator
            self._rect(slide, left + Inches(0.4), top + Inches(2.05),
                       col_w - Inches(0.8), Inches(0.01),
                       self.colors.border)

            # Context
            if c.context:
                self._text(slide, c.context,
                           left + Inches(0.15), top + Inches(2.2),
                           col_w - Inches(0.3), Inches(0.6),
                           size=9, color=self.colors.rgb("subtitle_text"),
                           align=PP_ALIGN.CENTER)

    # ── BAR CHART ────────────────────────────────────────────────────────

    def _bar_chart(self, slide, s: Slide, idx: int, total: int):
        self._add_title(slide, s.title)
        self._add_footer(slide, idx + 1, total, s.source)

        if not s.bars:
            return

        max_val = max(b.value for b in s.bars) or 1
        label_w = Inches(2.2)
        chart_left = self.grid.margin_left + label_w + Inches(0.15)
        annot_w = Inches(3.5)  # reserve space for annotations
        chart_w = self.grid.content_width - label_w - annot_w - Inches(0.3)
        bar_h = Inches(0.45)
        gap = Inches(0.12)
        y = self.grid.content_top + Inches(0.4)

        for bar in s.bars:
            # Label
            self._text(slide, bar.label,
                       self.grid.margin_left, y, label_w, bar_h,
                       size=self.typo.bullet_size, align=PP_ALIGN.RIGHT)

            # Bar
            bar_w = max(int(chart_w * (bar.value / max_val)), Inches(0.15))
            color = self.colors.primary if bar.highlight else ("#3A5068" if self.style.dark_mode else "#C0C0C0")
            self._rect(slide, chart_left, y + Inches(0.04), bar_w, bar_h - Inches(0.08), color)

            # Annotation
            if bar.annotation:
                self._text(slide, bar.annotation,
                           chart_left + bar_w + Inches(0.15), y,
                           Inches(3), bar_h,
                           size=self.typo.source_size + 1,
                           color=self.colors.rgb("subtitle_text"))

            y += bar_h + gap

    # ── PROCESS FLOW ─────────────────────────────────────────────────────

    def _process_flow(self, slide, s: Slide, idx: int, total: int):
        self._add_title(slide, s.title)
        self._add_footer(slide, idx + 1, total, s.source)

        # Prefer native rendering from flow_spec (vector text, pixel-perfect)
        if s.flow_spec:
            try:
                from sem_engine.schema import ProcessModel
                from sem_engine.renderers.pptx_native import PptxNativeFlowRenderer

                model = ProcessModel.from_dict(s.flow_spec)
                renderer = PptxNativeFlowRenderer(dark_mode=self.style.dark_mode)
                content_left = self.grid.margin_left / 914400  # EMU to inches
                content_top = (self.grid.content_top + Inches(0.2)) / 914400
                content_w = self.grid.content_width / 914400
                content_h = (self.grid.content_height - Inches(0.5)) / 914400
                renderer.render_on_slide(slide, model,
                                        left=content_left, top=content_top,
                                        width=content_w, height=content_h)
                return
            except Exception as e:
                # Fall through to image embedding - print error for debugging
                import traceback
                print(f"  Native flow render failed: {e}")
                traceback.print_exc()

        # Fallback: embed pre-rendered image
        if s.flow_image_path and Path(s.flow_image_path).exists():
            slide.shapes.add_picture(
                str(Path(s.flow_image_path)),
                self.grid.margin_left + Inches(0.3),
                self.grid.content_top + Inches(0.2),
                self.grid.content_width - Inches(0.6),
                self.grid.content_height - Inches(0.6)
            )
        else:
            self._text(slide, "[Process flow diagram]",
                       Inches(3), Inches(3.5), Inches(7), Inches(1),
                       size=16, color=self.colors.rgb("subtitle_text"),
                       align=PP_ALIGN.CENTER, italic=True)

    # ── FRAMEWORK / MATRIX ───────────────────────────────────────────────

    def _framework(self, slide, s: Slide, idx: int, total: int):
        self._add_title(slide, s.title)
        self._add_footer(slide, idx + 1, total, s.source)

        if not s.cells:
            return

        max_row = max(c.row for c in s.cells) + 1
        max_col = max(c.col for c in s.cells) + 1

        cell_gap = Inches(0.08)
        total_w = self.grid.content_width - Inches(1.5)
        total_h = self.grid.content_height - Inches(1.5)
        cell_w = int((total_w - cell_gap * (max_col - 1)) / max_col)
        cell_h = int((total_h - cell_gap * (max_row - 1)) / max_row)

        origin_x = self.grid.margin_left + Inches(1.2)
        origin_y = self.grid.content_top + Inches(0.4)

        # Color palette for cells
        cell_colors = [self.colors.secondary, self.colors.primary,
                       self.colors.accent_positive, self.colors.accent_neutral]

        for cell in s.cells:
            x = origin_x + cell.col * (cell_w + cell_gap)
            y = origin_y + cell.row * (cell_h + cell_gap)

            ci = (cell.row * max_col + cell.col) % len(cell_colors)
            shape = self._card(slide, x, y, cell_w, cell_h, cell_colors[ci])

            # Label
            self._text(slide, cell.label,
                       x + Inches(0.2), y + Inches(0.2),
                       cell_w - Inches(0.4), Inches(0.5),
                       size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                       align=PP_ALIGN.CENTER)

            if cell.description:
                self._text(slide, cell.description,
                           x + Inches(0.2), y + Inches(0.8),
                           cell_w - Inches(0.4), cell_h - Inches(1),
                           size=9, color=RGBColor(0xDD, 0xDD, 0xDD),
                           align=PP_ALIGN.CENTER)

        # Axis labels
        if s.x_axis:
            self._text(slide, s.x_axis,
                       origin_x, origin_y + total_h + Inches(0.15),
                       total_w, Inches(0.3),
                       size=10, bold=True, color=self.colors.rgb("subtitle_text"),
                       align=PP_ALIGN.CENTER)
        if s.y_axis:
            self._text(slide, s.y_axis,
                       self.grid.margin_left - Inches(0.1),
                       origin_y + int(total_h / 2) - Inches(0.2),
                       Inches(1.2), Inches(0.4),
                       size=10, bold=True, color=self.colors.rgb("subtitle_text"),
                       align=PP_ALIGN.CENTER)

    # ── TIMELINE ─────────────────────────────────────────────────────────

    def _timeline(self, slide, s: Slide, idx: int, total: int):
        self._add_title(slide, s.title)
        self._add_footer(slide, idx + 1, total, s.source)

        if not s.phases:
            return

        n = len(s.phases)
        gap = Inches(0.12)
        phase_w = self.grid.col_width(n, gap)
        top = self.grid.content_top + Inches(0.3)
        bar_h = Inches(0.5)

        for pi, ph in enumerate(s.phases):
            left = self.grid.col_left(pi, n, gap)
            color = ph.color or self.colors.primary

            # Phase bar
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, phase_w, bar_h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(color)
            shape.line.fill.background()
            shape.adjustments[0] = 0.15

            # Phase label in bar
            self._text(slide, ph.label,
                       left + Inches(0.1), top + Inches(0.07),
                       phase_w - Inches(0.2), Inches(0.35),
                       size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                       align=PP_ALIGN.CENTER)

            # Duration below bar
            if ph.duration:
                self._text(slide, ph.duration,
                           left, top + bar_h + Inches(0.08), phase_w, Inches(0.25),
                           size=9, color=self.colors.rgb("subtitle_text"),
                           align=PP_ALIGN.CENTER)

            # Items in card below
            if ph.items:
                items_top = top + bar_h + Inches(0.4)
                card_h = Inches(0.32) * len(ph.items) + Inches(0.3)
                self._card(slide, left, items_top, phase_w, card_h)

                y = items_top + Inches(0.15)
                for item in ph.items:
                    self._text(slide, item,
                               left + Inches(0.15), y, phase_w - Inches(0.3), Inches(0.28),
                               size=self.typo.bullet_size - 1)
                    y += Inches(0.3)

    # ── QUOTE ────────────────────────────────────────────────────────────

    def _quote(self, slide, s: Slide, idx: int, total: int):
        self._add_title(slide, s.title)
        self._add_footer(slide, idx + 1, total, s.source)

        quote = s.quote_text or s.body
        # Quote card
        card_top = Inches(2.0)
        self._card(slide, Inches(1.5), card_top, Inches(10), Inches(3.5))

        self._text(slide, f'"{quote}"',
                   Inches(2), Inches(2.5), Inches(9), Inches(2),
                   size=self.typo.quote_size, italic=True)

        if s.quote_attribution:
            self._text(slide, f"-- {s.quote_attribution}",
                       Inches(5), Inches(4.8), Inches(6), Inches(0.5),
                       size=13, color=self.colors.rgb("subtitle_text"),
                       align=PP_ALIGN.RIGHT)

    # ── CLOSING ──────────────────────────────────────────────────────────

    def _closing(self, slide, s: Slide, idx: int, total: int):
        # Keep same bg as rest of deck (dark stays dark)
        title = s.title or "Thank you"

        # Accent line
        self._rect(slide, Inches(1.2), Inches(2.5), Inches(2.5), Inches(0.04), self.colors.primary)

        self._text(slide, title,
                   Inches(1.2), Inches(2.7), Inches(10), Inches(1.2),
                   size=34, bold=True, color=self.colors.rgb("title_text"))

        if s.subtitle:
            self._text(slide, s.subtitle,
                       Inches(1.2), Inches(4.2), Inches(8), Inches(0.8),
                       size=14, color=self.colors.rgb("subtitle_text"))

        if s.bullets:
            y = Inches(5.2)
            for b in s.bullets:
                text = f"{b.lead}: {b.detail}" if b.detail else b.lead
                self._text(slide, text,
                           Inches(1.2), y, Inches(8), Inches(0.35),
                           size=11, color=self.colors.rgb("footer_text"))
                y += Inches(0.38)
