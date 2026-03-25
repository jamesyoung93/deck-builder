"""
Native PowerPoint diagram renderer.

Uses Graphviz ONLY for layout computation (node positions),
then draws everything as native python-pptx shapes:
- Rounded rectangles for nodes
- Lines/connectors for edges
- TextBoxes with Calibri font (identical to slide text)
- Rectangles for group containers

Result: text is vector, pixel-perfect, same sharpness as all other slide text.
"""

from __future__ import annotations

import json
import math
import subprocess
import shutil
from pathlib import Path

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from ..schema import ProcessModel, Node, Edge, Group, NodeType, EdgeType, ViewMode


# ── Color palettes ───────────────────────────────────────────────────────

GROUP_DARK = [
    {"fill": "#0E3854", "border": "#00A0DC", "label": "#00A0DC"},
    {"fill": "#0E3828", "border": "#00B894", "label": "#00B894"},
    {"fill": "#1F0E38", "border": "#9B59B6", "label": "#9B59B6"},
    {"fill": "#381E0E", "border": "#E67E22", "label": "#E67E22"},
    {"fill": "#380E0E", "border": "#E74C3C", "label": "#E74C3C"},
]

GROUP_LIGHT = [
    {"fill": "#E8F4FD", "border": "#2196F3", "label": "#1565C0"},
    {"fill": "#E8F8F0", "border": "#26A69A", "label": "#00796B"},
    {"fill": "#F0E8F8", "border": "#9C27B0", "label": "#7B1FA2"},
    {"fill": "#FDF0E8", "border": "#FF9800", "label": "#E65100"},
    {"fill": "#FDE8E8", "border": "#EF5350", "label": "#C62828"},
]

NODE_COLORS = {
    NodeType.PROCESS:    "#1565C0",
    NodeType.SUBPROCESS: "#00838F",
    NodeType.DATA:       "#6A1B9A",
    NodeType.GATE:       "#EF6C00",
    NodeType.START:      "#2E7D32",
    NodeType.END:        "#C62828",
    NodeType.EXTERNAL:   "#AD1457",
    NodeType.HUMAN:      "#00695C",
    NodeType.MONITOR:    "#7B1FA2",
    NodeType.DECISION:   "#F57F17",
    NodeType.LATENT:     "#546E7A",
    NodeType.VARIABLE:   "#1976D2",
}

EDGE_COLORS = {
    EdgeType.FLOW:         "#64B5F6",
    EdgeType.CAUSAL:       "#4FC3F7",
    EdgeType.FEEDBACK:     "#EF5350",
    EdgeType.CONDITIONAL:  "#FFA726",
    EdgeType.DATA_FLOW:    "#66BB6A",
    EdgeType.ASSOCIATION:  "#90A4AE",
    EdgeType.BIDIRECTIONAL:"#64B5F6",
    EdgeType.MODERATION:   "#CE93D8",
    EdgeType.MEDIATION:    "#4DB6AC",
}


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _wrap(label: str, max_len: int = 22) -> str:
    if len(label) <= max_len:
        return label
    words = label.split()
    lines, cur = [], ""
    for w in words:
        if cur and len(cur) + len(w) + 1 > max_len:
            lines.append(cur)
            cur = w
        else:
            cur = f"{cur} {w}" if cur else w
    if cur:
        lines.append(cur)
    return "\n".join(lines)


class PptxNativeFlowRenderer:
    """
    Renders a ProcessModel directly onto a PowerPoint slide as native shapes.

    Usage:
        renderer = PptxNativeFlowRenderer(dark_mode=True)
        renderer.render_on_slide(slide, model, content_area)
    """

    def __init__(self, dark_mode: bool = True,
                 node_w: float = 1.7, node_h: float = 0.6,
                 font_size: int = 10, label_font_size: int = 9,
                 group_label_size: int = 10):
        self.dark_mode = dark_mode
        self.node_w = node_w   # inches
        self.node_h = node_h
        self.font_size = font_size
        self.label_font_size = label_font_size
        self.group_label_size = group_label_size
        self.groups = GROUP_DARK if dark_mode else GROUP_LIGHT

    def render_on_slide(self, slide, model: ProcessModel,
                       left: float = 0.6, top: float = 1.5,
                       width: float = 12.0, height: float = 5.2):
        """
        Draw the process model directly on a pptx slide.

        Args:
            slide: pptx slide object
            model: ProcessModel to render
            left, top: content area origin in inches
            width, height: content area size in inches
        """
        # Get layout from Graphviz
        positions = self._compute_layout(model)
        if not positions:
            return

        # Scale positions to fit content area
        scaled = self._scale_positions(positions, left, top, width, height)

        # Draw groups first (background)
        self._draw_groups(slide, model, scaled)

        # Draw edges
        self._draw_edges(slide, model, scaled)

        # Draw nodes on top
        self._draw_nodes(slide, model, scaled)

    def _compute_layout(self, model: ProcessModel) -> dict[str, tuple[float, float]]:
        """Use Graphviz to compute node positions."""
        dot = self._make_dot(model)
        try:
            proc = subprocess.run(
                ["dot", "-Tjson"], input=dot,
                capture_output=True, text=True, timeout=15
            )
            if proc.returncode != 0:
                return {}
            data = json.loads(proc.stdout)
            return self._extract_positions(data)
        except Exception:
            return self._fallback_layout(model)

    def _make_dot(self, model: ProcessModel) -> str:
        lines = ['digraph layout {',
                 '  rankdir=LR;',
                 '  nodesep=1.0; ranksep=1.8;',
                 '  node [width=2.0, height=0.7, fixedsize=true];']

        for group in model.groups:
            nodes = model.get_nodes_in_group(group.id)
            if nodes:
                lines.append(f'  subgraph "cluster_{group.id}" {{')
                lines.append(f'    margin=25;')
                for n in nodes:
                    lines.append(f'    "{n.id}";')
                lines.append('  }')

        grouped = {n.id for g in model.groups for n in model.get_nodes_in_group(g.id)}
        for n in model.nodes:
            if n.id not in grouped:
                lines.append(f'  "{n.id}";')

        for e in model.edges:
            c = "false" if e.edge_type == EdgeType.FEEDBACK else "true"
            lines.append(f'  "{e.source}" -> "{e.target}" [constraint={c}];')

        lines.append('}')
        return '\n'.join(lines)

    def _extract_positions(self, data: dict) -> dict[str, tuple[float, float]]:
        positions = {}
        bb = data.get("bb", "0,0,100,100")
        parts = [float(x) for x in bb.split(",")]
        gv_h = parts[3] - parts[1]

        def extract(obj):
            if not isinstance(obj, dict):
                return
            name = obj.get("name", "")
            pos = obj.get("pos", "")
            # Node positions are simple "x,y" - edge positions have "e," prefix or spline data
            if pos and name and "," in pos and not pos.startswith("e,") and pos.count(",") == 1:
                px, py = pos.split(",")
                x = float(px) - parts[0]
                y = gv_h - (float(py) - parts[1])
                positions[name] = (x, y)
            # Recurse into sub-objects (clusters contain nodes)
            for sub in obj.get("objects", []):
                extract(sub)

        extract(data)
        return positions

    def _scale_positions(self, positions: dict, left: float, top: float,
                        width: float, height: float) -> dict[str, tuple[float, float]]:
        if not positions:
            return {}

        xs = [p[0] for p in positions.values()]
        ys = [p[1] for p in positions.values()]
        min_x, max_x = min(xs), max(xs)
        min_y, max_y = min(ys), max(ys)

        range_x = max_x - min_x or 1
        range_y = max_y - min_y or 1

        # Padding: room for node dimensions + group labels
        pad_x = self.node_w
        pad_y = self.node_h + 0.3  # extra for group header text

        avail_w = width - 2 * pad_x
        avail_h = height - 2 * pad_y

        # Scale uniformly to fit, then center in the content area
        scale_x = avail_w / range_x if range_x > 0 else 1
        scale_y = avail_h / range_y if range_y > 0 else 1

        # Use the full available space (stretch to fill)
        # If layout is very flat (all same y), center vertically
        if range_y < 10:  # Essentially a single row
            scale_y = 1
            y_offset = top + height / 2  # Center vertically
        else:
            y_offset = top + pad_y

        scaled = {}
        for name, (x, y) in positions.items():
            if range_y < 10:
                sx = left + pad_x + (x - min_x) * scale_x
                sy = y_offset
            else:
                sx = left + pad_x + (x - min_x) * scale_x
                sy = y_offset + (y - min_y) * scale_y
            scaled[name] = (sx, sy)

        return scaled

    def _fallback_layout(self, model: ProcessModel) -> dict[str, tuple[float, float]]:
        positions = {}
        n = len(model.nodes)
        cols = max(int(math.sqrt(n * 2)), 3)
        for i, node in enumerate(model.nodes):
            col = i % cols
            row = i // cols
            positions[node.id] = (col * 150 + 75, row * 100 + 50)
        return positions

    def _draw_groups(self, slide, model: ProcessModel, positions: dict):
        for gi, group in enumerate(model.groups):
            nodes = model.get_nodes_in_group(group.id)
            node_pos = [positions[n.id] for n in nodes if n.id in positions]
            if not node_pos:
                continue

            palette = self.groups[gi % len(self.groups)]

            margin_x = 0.25
            margin_top = 0.4   # room for label
            margin_bottom = 0.2
            min_x = min(p[0] for p in node_pos) - self.node_w / 2 - margin_x
            min_y = min(p[1] for p in node_pos) - self.node_h / 2 - margin_top
            max_x = max(p[0] for p in node_pos) + self.node_w / 2 + margin_x
            max_y = max(p[1] for p in node_pos) + self.node_h / 2 + margin_bottom

            gw = max_x - min_x
            gh = max_y - min_y

            # Group rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(min_x), Inches(min_y), Inches(gw), Inches(gh)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(palette["fill"])
            shape.line.color.rgb = _rgb(palette["border"])
            shape.line.width = Pt(1.5)
            try:
                shape.adjustments[0] = 0.03
            except Exception:
                pass

            # Accent bar at top of group
            bar_h = 0.04
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(min_x), Inches(min_y), Inches(gw), Inches(bar_h)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = _rgb(palette["border"])
            bar.line.fill.background()
            try:
                bar.adjustments[0] = 0.5
            except Exception:
                pass

            # Group label
            lbl = slide.shapes.add_textbox(
                Inches(min_x + 0.15), Inches(min_y + bar_h + 0.03),
                Inches(gw - 0.3), Inches(0.25)
            )
            p = lbl.text_frame.paragraphs[0]
            p.text = group.label
            p.font.size = Pt(self.group_label_size)
            p.font.bold = True
            p.font.name = "Calibri"
            p.font.color.rgb = _rgb(palette["label"])

    def _draw_nodes(self, slide, model: ProcessModel, positions: dict):
        for node in model.nodes:
            if node.id not in positions:
                continue

            cx, cy = positions[node.id]
            x = cx - self.node_w / 2
            y = cy - self.node_h / 2
            fill = NODE_COLORS.get(node.node_type, "#1565C0")

            # Node shape — rounded rectangle (pill)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(self.node_w), Inches(self.node_h)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
            shape.line.fill.background()
            try:
                shape.adjustments[0] = 0.35  # Very rounded = pill shape
            except Exception:
                pass

            # Node label — native PowerPoint text
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = _wrap(node.label, 20)
            p.font.size = Pt(self.font_size)
            p.font.bold = True
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

    def _draw_edges(self, slide, model: ProcessModel, positions: dict):
        for edge in model.edges:
            if edge.source not in positions or edge.target not in positions:
                continue

            sx, sy = positions[edge.source]
            tx, ty = positions[edge.target]
            color = EDGE_COLORS.get(edge.edge_type, "#64B5F6")

            # Adjust to node borders
            dx = tx - sx
            dy = ty - sy
            dist = math.sqrt(dx*dx + dy*dy) or 1

            # Start/end at node edges
            sx2 = sx + (dx / dist) * self.node_w * 0.55
            sy2 = sy + (dy / dist) * self.node_h * 0.55
            tx2 = tx - (dx / dist) * self.node_w * 0.55
            ty2 = ty - (dy / dist) * self.node_h * 0.55

            # Draw as a thin rectangle (line)
            # Calculate line geometry
            line_len = math.sqrt((tx2-sx2)**2 + (ty2-sy2)**2)
            if line_len < 0.05:
                continue

            # Use a simple line shape
            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR.STRAIGHT
                Inches(sx2), Inches(sy2),
                Inches(tx2), Inches(ty2)
            )
            connector.line.color.rgb = _rgb(color)
            connector.line.width = Pt(1.5)

            if edge.edge_type in (EdgeType.FEEDBACK, EdgeType.CONDITIONAL):
                connector.line.dash_style = 2  # dash

            # Arrowhead: small triangle at target end
            angle = math.atan2(ty2 - sy2, tx2 - sx2)
            arr_size = 0.1
            arr_shape = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                Inches(tx2 - arr_size/2), Inches(ty2 - arr_size/2),
                Inches(arr_size), Inches(arr_size)
            )
            arr_shape.fill.solid()
            arr_shape.fill.fore_color.rgb = _rgb(color)
            arr_shape.line.fill.background()
            arr_shape.rotation = math.degrees(angle) + 90

            # Edge label with background pill for readability
            if edge.label:
                mx = (sx2 + tx2) / 2
                my = (sy2 + ty2) / 2 - 0.15
                lbl_w = max(len(edge.label) * 0.07 + 0.25, 0.9)
                lbl_h = 0.25

                # Background pill
                bg_color = "#0D1B2A" if self.dark_mode else "#FFFFFF"
                pill = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(mx - lbl_w/2), Inches(my),
                    Inches(lbl_w), Inches(lbl_h)
                )
                pill.fill.solid()
                pill.fill.fore_color.rgb = _rgb(bg_color)
                pill.line.color.rgb = _rgb(color)
                pill.line.width = Pt(0.75)
                try:
                    pill.adjustments[0] = 0.4
                except Exception:
                    pass

                # Label text
                lbl = slide.shapes.add_textbox(
                    Inches(mx - lbl_w/2), Inches(my + 0.02),
                    Inches(lbl_w), Inches(lbl_h)
                )
                p = lbl.text_frame.paragraphs[0]
                p.text = edge.label
                p.font.size = Pt(self.label_font_size)
                p.font.bold = True
                p.font.name = "Calibri"
                p.font.color.rgb = _rgb(color)
                p.alignment = PP_ALIGN.CENTER
